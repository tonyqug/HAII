from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas


def make_pdf(path: Path) -> None:
    c = canvas.Canvas(str(path), pagesize=letter)
    width, height = letter
    slides = [
        ("Biology 101 - Cell Structure", ["Nucleus stores DNA", "Mitochondria generate ATP", "Cell membrane regulates transport"]),
        ("Biology 101 - Photosynthesis", ["Occurs in chloroplasts", "Light reactions make ATP and NADPH", "Calvin cycle fixes carbon"]),
    ]
    for title, bullets in slides:
        c.setFont("Helvetica-Bold", 24)
        c.drawString(72, height - 72, title)
        c.setFont("Helvetica", 15)
        y = height - 120
        for bullet in bullets:
            c.drawString(96, y, f"• {bullet}")
            y -= 28
        c.showPage()
    c.save()


def make_pptx(path: Path) -> None:
    prs = Presentation()
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "History 201 - Causes of World War I"
    slide.placeholders[1].text = "Alliance system\nMilitarism\nImperial competition\nNationalism"

    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "History 201 - Immediate Trigger"
    slide.placeholders[1].text = "Assassination of Archduke Franz Ferdinand\nJuly Crisis escalated quickly"
    prs.save(str(path))


def make_notes(path: Path) -> None:
    path.write_text(
        "# Sample Notes\n\n"
        "The mitochondrion is the site of cellular respiration.\n\n"
        "## Exam reminder\n"
        "Understand how ATP is produced and how chloroplasts differ from mitochondria.\n",
        encoding="utf-8",
    )


def main() -> None:
    out_dir = Path(__file__).resolve().parent
    make_pdf(out_dir / "sample_slides.pdf")
    make_pptx(out_dir / "sample_deck.pptx")
    make_notes(out_dir / "sample_notes.md")


if __name__ == "__main__":
    main()
