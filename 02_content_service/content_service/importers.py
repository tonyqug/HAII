from __future__ import annotations

import json
import traceback
from pathlib import Path
from typing import Any, Optional

import fitz
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .config import Settings
from .rendering import convert_pptx_to_pdf, relative_path, render_pdf_page_to_png, render_text_image
from .repository import Repository
from .utils import (
    estimate_tokens,
    extract_title_guess,
    now_iso,
    safe_filename,
    sha1_text,
    split_text_into_units,
    summarize_quality,
)


def material_base_dir(settings: Settings, workspace_id: str, material_id: str) -> Path:
    return settings.storage_dir / "workspaces" / workspace_id / "materials" / material_id


def source_file_path(settings: Settings, workspace_id: str, material_id: str, original_name: str) -> Path:
    return material_base_dir(settings, workspace_id, material_id) / "source" / safe_filename(original_name)


def derived_dir(settings: Settings, workspace_id: str, material_id: str) -> Path:
    return material_base_dir(settings, workspace_id, material_id) / "derived"


def previews_dir(settings: Settings, workspace_id: str, material_id: str) -> Path:
    return material_base_dir(settings, workspace_id, material_id) / "previews"


def write_normalized_outputs(settings: Settings, workspace_id: str, material_id: str, material_meta: dict[str, Any], slides: list[dict[str, Any]]) -> None:
    ddir = derived_dir(settings, workspace_id, material_id)
    ddir.mkdir(parents=True, exist_ok=True)
    manifest = {
        "material_id": material_id,
        "workspace_id": workspace_id,
        "exported_at": now_iso(),
        "material": material_meta,
        "slides": [
            {
                "slide_id": slide["slide_id"],
                "slide_number": slide["slide_number"],
                "title_guess": slide.get("title_guess"),
                "extracted_text": slide.get("extracted_text", ""),
                "extraction_quality": slide["extraction_quality"],
                "quality_notes": slide.get("quality_notes"),
                "has_text": bool(slide.get("has_text")),
                "preview_relpath": slide.get("preview_relpath"),
            }
            for slide in slides
        ],
    }
    (ddir / "manifest.json").write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")
    (ddir / "slides.json").write_text(json.dumps(manifest["slides"], ensure_ascii=False, indent=2), encoding="utf-8")


def _slide_id(material_id: str, slide_number: int) -> str:
    return f"{material_id}:slide:{slide_number}"


def _extract_shape_text(shape: Any) -> list[str]:
    texts: list[str] = []
    if hasattr(shape, "text") and shape.text:
        text = str(shape.text).strip()
        if text:
            texts.append(text)
    if getattr(shape, "has_text_frame", False) and getattr(shape, "text_frame", None) is not None:
        frame_text = shape.text_frame.text.strip()
        if frame_text and frame_text not in texts:
            texts.append(frame_text)
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                value = cell.text.strip()
                if value:
                    texts.append(value)
    if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
        for subshape in shape.shapes:
            texts.extend(_extract_shape_text(subshape))
    return texts


def process_material_import(settings: Settings, repo: Repository, material_id: str, job_id: str) -> None:
    material = repo.get_material(material_id)
    if material is None:
        repo.update_job(
            job_id,
            status="failed",
            progress=100,
            stage="missing_material",
            message="The reserved material record was missing.",
            error_code="material_not_found",
            error_message="Material record disappeared before processing started.",
            error_retryable=0,
        )
        return

    try:
        repo.update_job(job_id, status="running", progress=2, stage="starting", message="Starting material import.")
        repo.update_material(material_id, processing_status="running")

        kind = material["kind"]
        workspace_id = material["workspace_id"]
        base_dir = material_base_dir(settings, workspace_id, material_id)
        base_dir.mkdir(parents=True, exist_ok=True)
        source_path = settings.local_data_dir / material["source_relpath"] if material.get("source_relpath") else None
        if source_path is None or not source_path.exists():
            raise RuntimeError("Source payload was not stored before processing began.")

        if kind == "pdf":
            meta, slides = _process_pdf(settings, repo, material, job_id, source_path)
        elif kind == "pptx":
            meta, slides = _process_pptx(settings, repo, material, job_id, source_path)
        elif kind in {"pasted_text", "text"}:
            meta, slides = _process_text_like(settings, repo, material, job_id, source_path)
        else:
            raise RuntimeError(f"Unsupported material kind: {kind}")

        write_normalized_outputs(settings, workspace_id, material_id, meta, slides)
        repo.replace_slides(material_id, slides)
        repo.update_material(
            material_id,
            processing_status="ready",
            page_count=meta["page_count"],
            slide_count=meta["slide_count"],
            ready_for_retrieval=1 if meta["ready_for_retrieval"] else 0,
            previews_available=1 if meta["previews_available"] else 0,
            quality_overall=meta["quality_overall"],
            quality_notes=meta.get("quality_notes"),
            extraction_notes=meta.get("extraction_notes"),
            converted_pdf_relpath=meta.get("converted_pdf_relpath"),
            failure_code=None,
            failure_message=None,
        )
        repo.update_job(
            job_id,
            status="succeeded",
            progress=100,
            stage="completed",
            message="Material import completed successfully.",
            result_type="material",
            result_id=material_id,
            error_code=None,
            error_message=None,
            error_retryable=None,
        )
    except Exception as exc:  # noqa: BLE001 - failure should be surfaced as job state
        failure_message = str(exc).strip() or exc.__class__.__name__
        trace = traceback.format_exc(limit=8)
        repo.update_material(
            material_id,
            processing_status="failed",
            ready_for_retrieval=0,
            previews_available=0,
            quality_overall="low",
            quality_notes="Import failed.",
            extraction_notes=trace[:4000],
            failure_code="import_failed",
            failure_message=failure_message,
        )
        repo.update_job(
            job_id,
            status="failed",
            progress=100,
            stage="failed",
            message=f"Material import failed: {failure_message}",
            result_type="material",
            result_id=material_id,
            error_code="import_failed",
            error_message=failure_message,
            error_retryable=1,
        )


def _process_pdf(settings: Settings, repo: Repository, material: dict[str, Any], job_id: str, pdf_path: Path) -> tuple[dict[str, Any], list[dict[str, Any]]]:
    repo.update_job(job_id, progress=8, stage="opening_pdf", message="Opening PDF and reading page count.")
    slides: list[dict[str, Any]] = []
    preview_count = 0
    extracted_count = 0
    with fitz.open(pdf_path) as doc:
        total = doc.page_count
        if total <= 0:
            raise RuntimeError("The uploaded PDF contained no pages.")
        for index, page in enumerate(doc, start=1):
            text = page.get_text("text").strip()
            if not text:
                blocks = page.get_text("blocks")
                if blocks:
                    text = "\n".join(str(block[4]).strip() for block in blocks if len(block) >= 5 and str(block[4]).strip())
            title_guess = extract_title_guess(text)
            preview_path = previews_dir(settings, material["workspace_id"], material["material_id"]) / f"slide-{index:04d}.png"
            preview_relpath: Optional[str] = None
            quality = "low"
            quality_notes = None
            try:
                render_pdf_page_to_png(pdf_path, index - 1, preview_path)
                preview_count += 1
                preview_relpath = relative_path(settings.local_data_dir, preview_path)
            except Exception as exc:  # noqa: BLE001
                quality_notes = f"Preview generation failed: {exc}"
            if len(text) >= 60:
                quality = "high" if preview_relpath else "medium"
                extracted_count += 1
            elif len(text) >= 8:
                quality = "medium" if preview_relpath else "low"
                extracted_count += 1
                quality_notes = quality_notes or "Text extraction was partial or noisy."
            else:
                quality = "low"
                quality_notes = quality_notes or "No readable text was extracted from this page."
            slide = {
                "slide_id": _slide_id(material["material_id"], index),
                "slide_number": index,
                "title_guess": title_guess,
                "extracted_text": text,
                "extraction_quality": quality,
                "quality_notes": quality_notes,
                "has_text": bool(text.strip()),
                "preview_relpath": preview_relpath,
                "token_count": estimate_tokens(text),
                "text_checksum": sha1_text(text),
            }
            slides.append(slide)
            progress = min(94, 8 + int((index / total) * 84))
            repo.update_job(job_id, progress=progress, stage="processing_pdf_pages", message=f"Processed PDF page {index} of {total}.")

    qualities = [slide["extraction_quality"] for slide in slides]
    previews_available = preview_count > 0
    overall, note = summarize_quality(qualities, previews_available, extracted_count / max(len(slides), 1))
    meta = {
        "page_count": len(slides),
        "slide_count": len(slides),
        "ready_for_retrieval": extracted_count > 0,
        "previews_available": previews_available,
        "quality_overall": overall,
        "quality_notes": note,
        "extraction_notes": "PDF text was extracted with PyMuPDF page parsing. No OCR text was fabricated.",
        "converted_pdf_relpath": relative_path(settings.local_data_dir, pdf_path),
    }
    return meta, slides


def _process_pptx(settings: Settings, repo: Repository, material: dict[str, Any], job_id: str, pptx_path: Path) -> tuple[dict[str, Any], list[dict[str, Any]]]:
    repo.update_job(job_id, progress=8, stage="reading_pptx", message="Reading PPTX slide structure and text.")
    presentation = Presentation(str(pptx_path))
    total_slides = len(presentation.slides)
    if total_slides <= 0:
        raise RuntimeError("The uploaded PPTX contained no slides.")

    slide_texts: list[dict[str, Any]] = []
    extracted_count = 0
    for idx, slide in enumerate(presentation.slides, start=1):
        pieces: list[str] = []
        for shape in slide.shapes:
            pieces.extend(_extract_shape_text(shape))
        text = "\n".join(piece for piece in pieces if piece.strip()).strip()
        title_guess = None
        if getattr(slide.shapes, "title", None) is not None and slide.shapes.title is not None:
            try:
                title_text = (slide.shapes.title.text or "").strip()
                if title_text:
                    title_guess = title_text[:160]
            except Exception:  # noqa: BLE001
                title_guess = None
        title_guess = title_guess or extract_title_guess(text)
        if text:
            extracted_count += 1
        slide_texts.append(
            {
                "slide_id": _slide_id(material["material_id"], idx),
                "slide_number": idx,
                "title_guess": title_guess,
                "extracted_text": text,
                "extraction_quality": "high" if len(text) >= 20 else ("medium" if text else "low"),
                "quality_notes": None if text else "No readable text was found in PowerPoint shapes on this slide.",
                "has_text": bool(text),
                "preview_relpath": None,
                "token_count": estimate_tokens(text),
                "text_checksum": sha1_text(text),
            }
        )
        progress = min(45, 8 + int((idx / total_slides) * 37))
        repo.update_job(job_id, progress=progress, stage="extracting_pptx_text", message=f"Extracted text from PPTX slide {idx} of {total_slides}.")

    converted_pdf_relpath: Optional[str] = None
    preview_count = 0
    preview_failure: Optional[str] = None
    pdf_path: Optional[Path] = None
    try:
        repo.update_job(job_id, progress=50, stage="rendering_pptx", message="Converting PPTX to PDF for stable slide previews.")
        pdf_dir = material_base_dir(settings, material["workspace_id"], material["material_id"]) / "rendered"
        pdf_path = convert_pptx_to_pdf(pptx_path, pdf_dir, soffice_bin=settings.libreoffice_bin)
        converted_pdf_relpath = relative_path(settings.local_data_dir, pdf_path)
        with fitz.open(pdf_path) as doc:
            for slide in slide_texts:
                idx = slide["slide_number"]
                preview_path = previews_dir(settings, material["workspace_id"], material["material_id"]) / f"slide-{idx:04d}.png"
                if idx - 1 < doc.page_count:
                    render_pdf_page_to_png(pdf_path, idx - 1, preview_path)
                    slide["preview_relpath"] = relative_path(settings.local_data_dir, preview_path)
                    preview_count += 1
                else:
                    render_text_image(
                        slide["extracted_text"] or "Preview was unavailable for this slide.",
                        preview_path,
                        title=slide.get("title_guess") or f"Slide {idx}",
                        footer="PowerPoint preview fallback",
                    )
                    slide["preview_relpath"] = relative_path(settings.local_data_dir, preview_path)
                progress = min(92, 50 + int((idx / total_slides) * 42))
                repo.update_job(job_id, progress=progress, stage="rendering_pptx", message=f"Rendered PPTX preview {idx} of {total_slides}.")
    except Exception as exc:  # noqa: BLE001
        preview_failure = str(exc)
        repo.update_job(job_id, progress=55, stage="rendering_pptx_fallback", message="PPTX conversion failed; generating text-based placeholder previews.")
        for slide in slide_texts:
            idx = slide["slide_number"]
            preview_path = previews_dir(settings, material["workspace_id"], material["material_id"]) / f"slide-{idx:04d}.png"
            render_text_image(
                slide["extracted_text"] or "No readable slide text was extracted.",
                preview_path,
                title=slide.get("title_guess") or f"Slide {idx}",
                footer="Text-based fallback preview",
            )
            slide["preview_relpath"] = relative_path(settings.local_data_dir, preview_path)
            if slide["extraction_quality"] == "high":
                slide["extraction_quality"] = "medium"
            slide["quality_notes"] = slide.get("quality_notes") or f"Preview used a text fallback because PPTX conversion failed: {preview_failure}"
            progress = min(92, 55 + int((idx / total_slides) * 37))
            repo.update_job(job_id, progress=progress, stage="rendering_pptx_fallback", message=f"Built fallback preview {idx} of {total_slides}.")

    previews_available = any(slide.get("preview_relpath") for slide in slide_texts)
    qualities = [slide["extraction_quality"] for slide in slide_texts]
    overall, note = summarize_quality(qualities, previews_available, extracted_count / max(total_slides, 1))
    extraction_notes = "PPTX text was extracted from slide shapes via python-pptx. Previews were rendered through LibreOffice PDF conversion."
    if preview_failure:
        extraction_notes += f" Preview fallback was used because conversion failed: {preview_failure}"
    meta = {
        "page_count": total_slides,
        "slide_count": total_slides,
        "ready_for_retrieval": extracted_count > 0,
        "previews_available": previews_available,
        "quality_overall": overall,
        "quality_notes": note,
        "extraction_notes": extraction_notes,
        "converted_pdf_relpath": converted_pdf_relpath,
    }
    return meta, slide_texts


def _process_text_like(settings: Settings, repo: Repository, material: dict[str, Any], job_id: str, text_path: Path) -> tuple[dict[str, Any], list[dict[str, Any]]]:
    raw_text = text_path.read_text(encoding="utf-8", errors="ignore")
    repo.update_job(job_id, progress=8, stage="reading_text", message="Reading text and splitting it into retrievable units.")
    chunks = split_text_into_units(raw_text)
    slides: list[dict[str, Any]] = []
    non_empty = 0
    for idx, chunk in enumerate(chunks, start=1):
        preview_path = previews_dir(settings, material["workspace_id"], material["material_id"]) / f"slide-{idx:04d}.png"
        render_text_image(
            chunk or "No text was provided in this section.",
            preview_path,
            title=extract_title_guess(chunk) or material["title"] or f"Page {idx}",
            footer=f"{material['kind']} section {idx}",
        )
        if chunk.strip():
            non_empty += 1
        slides.append(
            {
                "slide_id": _slide_id(material["material_id"], idx),
                "slide_number": idx,
                "title_guess": extract_title_guess(chunk),
                "extracted_text": chunk,
                "extraction_quality": "high" if chunk.strip() else "low",
                "quality_notes": None if chunk.strip() else "This section was empty.",
                "has_text": bool(chunk.strip()),
                "preview_relpath": relative_path(settings.local_data_dir, preview_path),
                "token_count": estimate_tokens(chunk),
                "text_checksum": sha1_text(chunk),
            }
        )
        progress = min(95, 8 + int((idx / max(len(chunks), 1)) * 87))
        repo.update_job(job_id, progress=progress, stage="rendering_text", message=f"Prepared text section {idx} of {len(chunks)}.")
    qualities = [slide["extraction_quality"] for slide in slides]
    overall, note = summarize_quality(qualities, True, non_empty / max(len(slides), 1))
    meta = {
        "page_count": len(slides),
        "slide_count": len(slides),
        "ready_for_retrieval": non_empty > 0,
        "previews_available": True,
        "quality_overall": overall,
        "quality_notes": note,
        "extraction_notes": "User-supplied text was normalized into local sections without adding any generated wording.",
        "converted_pdf_relpath": None,
    }
    return meta, slides
